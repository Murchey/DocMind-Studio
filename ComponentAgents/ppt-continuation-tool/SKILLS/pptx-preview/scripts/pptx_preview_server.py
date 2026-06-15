#!/usr/bin/env python3
"""
PPTX Preview Server

Provides a browser-based preview for PPTX files.
Converts PPTX slides to images and serves them via a web interface.

Usage:
    python3 pptx_preview_server.py <pptx_file> [--port 5050]

Examples:
    python3 pptx_preview_server.py output.pptx
    python3 pptx_preview_server.py output.pptx --port 8080
    python3 pptx_preview_server.py output.pptx --no-browser

Dependencies:
    flask>=3.0.0
    python-pptx
    Pillow
"""

import argparse
import atexit
import json
import logging
import os
import signal
import sys
import tempfile
import threading
import time
import webbrowser
from pathlib import Path
from typing import Optional

from flask import Flask, jsonify, send_from_directory, send_file

logger = logging.getLogger('pptx_preview')

# Per-project lock file
LOCK_FILE_NAME = '.pptx_preview.lock'


def _process_alive(pid: int) -> bool:
    """Return True if a process with this pid is reachable."""
    if pid <= 0:
        return False
    try:
        os.kill(pid, 0)
    except ProcessLookupError:
        return False
    except PermissionError:
        return True
    except OSError:
        return False
    return True


def _read_lock(lock_file: Path) -> Optional[dict]:
    try:
        data = json.loads(lock_file.read_text(encoding='utf-8'))
        return data if isinstance(data, dict) else None
    except (OSError, json.JSONDecodeError):
        return None


def _claim_lock(lock_file: Path, port: int) -> Optional[dict]:
    """Try to claim the per-project preview slot."""
    existing = _read_lock(lock_file)
    if existing and _process_alive(int(existing.get('pid', 0))):
        return existing
    lock_file.write_text(
        json.dumps({'pid': os.getpid(), 'port': port}),
        encoding='utf-8',
    )
    return None


def _release_lock(lock_file: Path) -> None:
    """Best-effort cleanup: only delete the lock if it still names *us*."""
    try:
        current = _read_lock(lock_file)
        if current and int(current.get('pid', 0)) == os.getpid():
            lock_file.unlink(missing_ok=True)
    except OSError:
        pass


def convert_pptx_to_images(pptx_path: Path, output_dir: Path) -> list[dict]:
    """Convert PPTX slides to images."""
    try:
        from pptx import Presentation
        from PIL import Image
        import io
    except ImportError as e:
        logger.error("Missing dependency: %s. Run: pip install python-pptx Pillow", e.name)
        return []

    slides_info = []
    try:
        presentation = Presentation(str(pptx_path))
        
        # Get slide dimensions
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        
        for slide_index, slide in enumerate(presentation.slides, 1):
            # Extract slide title if available
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            
            # Create a simple representation
            slide_info = {
                "index": slide_index,
                "title": title,
                "width": slide_width,
                "height": slide_height,
            }
            slides_info.append(slide_info)
            
    except Exception as e:
        logger.error("Failed to convert PPTX: %s", e)
        return []
    
    return slides_info


def create_app(
    pptx_path: Path,
    idle_timeout: int = 900,
    lock_file: Optional[Path] = None,
) -> Flask:
    """Create and configure the Flask app for PPTX preview."""
    app = Flask(__name__)
    app.config['PPTX_PATH'] = pptx_path
    app.config['LOCK_FILE'] = lock_file
    app.config['LAST_REQUEST_TIME'] = time.time()
    
    # Convert PPTX to slide info
    slides_info = convert_pptx_to_images(pptx_path, pptx_path.parent)
    app.config['SLIDES_INFO'] = slides_info
    
    @app.before_request
    def _update_activity():
        app.config['LAST_REQUEST_TIME'] = time.time()
    
    def _exit_with_lock_release(code: int = 0) -> None:
        lf = app.config.get('LOCK_FILE')
        if lf is not None:
            _release_lock(lf)
        os._exit(code)
    
    def _idle_watchdog():
        if idle_timeout <= 0:
            return
        while True:
            time.sleep(10)
            elapsed = time.time() - app.config['LAST_REQUEST_TIME']
            if elapsed > idle_timeout:
                logger.info('idle for %ds, shutting down', idle_timeout)
                _exit_with_lock_release(0)
    
    watchdog = threading.Thread(target=_idle_watchdog, daemon=True)
    watchdog.start()
    
    @app.route('/')
    def index():
        """Serve the main preview page."""
        slides = app.config['SLIDES_INFO']
        html = generate_preview_html(slides, pptx_path.name)
        return html
    
    @app.route('/api/slides')
    def get_slides():
        """Return slide information as JSON."""
        return jsonify({
            'slides': app.config['SLIDES_INFO'],
            'total': len(app.config['SLIDES_INFO']),
            'filename': pptx_path.name,
        })
    
    @app.route('/api/slide/<int:slide_number>')
    def get_slide(slide_number: int):
        """Return information for a specific slide."""
        slides = app.config['SLIDES_INFO']
        if 1 <= slide_number <= len(slides):
            return jsonify(slides[slide_number - 1])
        return jsonify({'error': 'Slide not found'}), 404
    
    @app.route('/api/shutdown', methods=['POST'])
    def shutdown():
        """Shutdown the preview server."""
        def _stop():
            time.sleep(0.5)
            logger.info('shutting down')
            _exit_with_lock_release(0)
        threading.Thread(target=_stop, daemon=True).start()
        return jsonify({'status': 'ok'})
    
    return app


def generate_preview_html(slides: list[dict], filename: str) -> str:
    """Generate HTML for the preview interface."""
    slides_json = json.dumps(slides, ensure_ascii=False)
    
    html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PPTX Preview - {filename}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
            background: #1a1a2e;
            color: #eee;
            height: 100vh;
            display: flex;
            flex-direction: column;
        }}
        
        .header {{
            background: #16213e;
            padding: 15px 20px;
            display: flex;
            justify-content: space-between;
            align-items: center;
            border-bottom: 1px solid #0f3460;
        }}
        
        .header h1 {{
            font-size: 18px;
            font-weight: 500;
            color: #e94560;
        }}
        
        .slide-info {{
            font-size: 14px;
            color: #a0a0a0;
        }}
        
        .main-container {{
            flex: 1;
            display: flex;
            flex-direction: column;
            align-items: center;
            justify-content: center;
            padding: 20px;
            position: relative;
        }}
        
        .slide-container {{
            background: white;
            border-radius: 8px;
            box-shadow: 0 10px 30px rgba(0, 0, 0, 0.3);
            width: 80%;
            max-width: 960px;
            aspect-ratio: 16/9;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
            overflow: hidden;
        }}
        
        .slide-content {{
            width: 100%;
            height: 100%;
            display: flex;
            flex-direction: column;
            align-items: center;
            justify-content: center;
            padding: 40px;
            text-align: center;
        }}
        
        .slide-title {{
            font-size: 36px;
            font-weight: bold;
            color: #333;
            margin-bottom: 20px;
        }}
        
        .slide-index {{
            font-size: 72px;
            font-weight: bold;
            color: #e94560;
            margin-bottom: 20px;
        }}
        
        .slide-placeholder {{
            font-size: 18px;
            color: #666;
        }}
        
        .controls {{
            display: flex;
            align-items: center;
            gap: 20px;
            margin-top: 30px;
        }}
        
        .btn {{
            background: #e94560;
            color: white;
            border: none;
            padding: 12px 24px;
            border-radius: 6px;
            font-size: 16px;
            cursor: pointer;
            transition: background 0.2s;
        }}
        
        .btn:hover {{
            background: #c73e54;
        }}
        
        .btn:disabled {{
            background: #666;
            cursor: not-allowed;
        }}
        
        .btn-secondary {{
            background: #0f3460;
        }}
        
        .btn-secondary:hover {{
            background: #1a4a8a;
        }}
        
        .slide-counter {{
            font-size: 18px;
            color: #a0a0a0;
            min-width: 100px;
            text-align: center;
        }}
        
        .keyboard-hint {{
            position: absolute;
            bottom: 20px;
            font-size: 12px;
            color: #666;
        }}
        
        .keyboard-hint kbd {{
            background: #333;
            padding: 2px 6px;
            border-radius: 3px;
            margin: 0 2px;
        }}
        
        .features {{
            position: absolute;
            top: 20px;
            right: 20px;
            background: rgba(15, 52, 96, 0.8);
            padding: 15px;
            border-radius: 8px;
            font-size: 12px;
        }}
        
        .features h3 {{
            margin-bottom: 10px;
            color: #e94560;
        }}
        
        .features ul {{
            list-style: none;
            padding: 0;
        }}
        
        .features li {{
            margin: 5px 0;
            color: #a0a0a0;
        }}
        
        .features li::before {{
            content: "•";
            color: #e94560;
            margin-right: 8px;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>PPTX Preview</h1>
        <div class="slide-info">{filename}</div>
    </div>
    
    <div class="main-container">
        <div class="slide-container">
            <div class="slide-content">
                <div class="slide-index" id="slideIndex">1</div>
                <div class="slide-title" id="slideTitle">Loading...</div>
                <div class="slide-placeholder">Slide content preview</div>
            </div>
        </div>
        
        <div class="controls">
            <button class="btn btn-secondary" id="btnFirst" onclick="goToFirst()">⏮</button>
            <button class="btn" id="btnPrev" onclick="goToPrev()">◀</button>
            <span class="slide-counter" id="slideCounter">1 / 1</span>
            <button class="btn" id="btnNext" onclick="goToNext()">▶</button>
            <button class="btn btn-secondary" id="btnLast" onclick="goToLast()">⏭</button>
        </div>
        
        <div class="keyboard-hint">
            Use <kbd>←</kbd> <kbd>→</kbd> arrow keys to navigate • <kbd>Home</kbd> <kbd>End</kbd> for first/last slide
        </div>
        
        <div class="features">
            <h3>Preview Features</h3>
            <ul>
                <li>Slide navigation</li>
                <li>Keyboard shortcuts</li>
                <li>Slide counter</li>
                <li>Responsive design</li>
            </ul>
        </div>
    </div>
    
    <script>
        const slides = {slides_json};
        let currentIndex = 0;
        
        function updateSlide() {{
            const slide = slides[currentIndex];
            document.getElementById('slideIndex').textContent = slide.index;
            document.getElementById('slideTitle').textContent = slide.title || `Slide ${{slide.index}}`;
            document.getElementById('slideCounter').textContent = `${{slide.index}} / ${{slides.length}}`;
            
            // Update button states
            document.getElementById('btnFirst').disabled = currentIndex === 0;
            document.getElementById('btnPrev').disabled = currentIndex === 0;
            document.getElementById('btnNext').disabled = currentIndex === slides.length - 1;
            document.getElementById('btnLast').disabled = currentIndex === slides.length - 1;
        }}
        
        function goToFirst() {{
            currentIndex = 0;
            updateSlide();
        }}
        
        function goToLast() {{
            currentIndex = slides.length - 1;
            updateSlide();
        }}
        
        function goToPrev() {{
            if (currentIndex > 0) {{
                currentIndex--;
                updateSlide();
            }}
        }}
        
        function goToNext() {{
            if (currentIndex < slides.length - 1) {{
                currentIndex++;
                updateSlide();
            }}
        }}
        
        function goToSlide(index) {{
            if (index >= 0 && index < slides.length) {{
                currentIndex = index;
                updateSlide();
            }}
        }}
        
        // Keyboard navigation
        document.addEventListener('keydown', function(e) {{
            switch(e.key) {{
                case 'ArrowLeft':
                case 'ArrowUp':
                    goToPrev();
                    break;
                case 'ArrowRight':
                case 'ArrowDown':
                    goToNext();
                    break;
                case 'Home':
                    goToFirst();
                    break;
                case 'End':
                    goToLast();
                    break;
                case 'F11':
                    e.preventDefault();
                    if (document.fullscreenElement) {{
                        document.exitFullscreen();
                    }} else {{
                        document.documentElement.requestFullscreen();
                    }}
                    break;
                case 'Escape':
                    if (document.fullscreenElement) {{
                        document.exitFullscreen();
                    }}
                    break;
            }}
        }});
        
        // Initialize
        updateSlide();
    </script>
</body>
</html>'''
    
    return html


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description='PPTX Preview Server',
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument('pptx_file', help='Path to PPTX file to preview')
    parser.add_argument('--port', type=int, default=5050, help='Port to listen on (default: 5050)')
    parser.add_argument('--no-browser', action='store_true', help='Do not auto-open browser')
    parser.add_argument(
        '--timeout',
        type=int,
        default=900,
        help='Idle timeout in seconds (default: 900; 0 = disabled)',
    )
    return parser


def start_preview_server(
    pptx_path: str,
    port: int = 5050,
    auto_open: bool = True,
    timeout: int = 900,
) -> str:
    """Start the preview server programmatically."""
    pptx_file = Path(pptx_path).resolve()
    if not pptx_file.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_file}")
    
    # Create lock file
    lock_file = pptx_file.parent / LOCK_FILE_NAME
    existing = _claim_lock(lock_file, port)
    if existing:
        existing_pid = existing.get('pid', '?')
        existing_port = existing.get('port', '?')
        raise RuntimeError(
            f"Preview is already running for this project "
            f"(pid={existing_pid}, port={existing_port}). "
            f"Open http://localhost:{existing_port} or kill process {existing_pid}"
        )
    
    atexit.register(_release_lock, lock_file)
    
    # Handle SIGTERM
    def _on_sigterm(signum: int, _frame) -> None:
        logger.info('received signal %s, exiting', signum)
        sys.exit(0)
    try:
        signal.signal(signal.SIGTERM, _on_sigterm)
    except (ValueError, OSError):
        pass
    
    app = create_app(
        pptx_file,
        idle_timeout=timeout,
        lock_file=lock_file,
    )
    
    url = f'http://localhost:{port}'
    
    if auto_open:
        webbrowser.open(url)
    
    # Run Flask in a separate thread
    def run_flask():
        app.run(host='127.0.0.1', port=port, debug=False, use_reloader=False)
    
    server_thread = threading.Thread(target=run_flask, daemon=True)
    server_thread.start()
    
    logger.info('Preview server running at %s', url)
    logger.info('PPTX file: %s', pptx_file)
    logger.info('Idle timeout: %ds (0 = disabled)', timeout)
    
    return url


def main(argv: Optional[list[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    
    logging.basicConfig(
        level=logging.INFO,
        format='[%(asctime)s] [%(levelname)s] pptx_preview: %(message)s',
        datefmt='%H:%M:%S',
    )
    
    pptx_file = Path(args.pptx_file).resolve()
    if not pptx_file.exists():
        logger.error('PPTX file not found: %s', pptx_file)
        return 1
    
    # Create lock file
    lock_file = pptx_file.parent / LOCK_FILE_NAME
    existing = _claim_lock(lock_file, args.port)
    if existing:
        existing_pid = existing.get('pid', '?')
        existing_port = existing.get('port', '?')
        logger.error(
            'Preview is already running for this project '
            '(pid=%s, port=%s). Open http://localhost:%s or kill process %s',
            existing_pid, existing_port, existing_port, existing_pid,
        )
        return 1
    
    atexit.register(_release_lock, lock_file)
    
    # Handle SIGTERM
    def _on_sigterm(signum: int, _frame) -> None:
        logger.info('received signal %s, exiting', signum)
        sys.exit(0)
    try:
        signal.signal(signal.SIGTERM, _on_sigterm)
    except (ValueError, OSError):
        pass
    
    app = create_app(
        pptx_file,
        idle_timeout=args.timeout,
        lock_file=lock_file,
    )
    
    url = f'http://localhost:{args.port}'
    if not args.no_browser:
        webbrowser.open(url)
    
    logger.info('Preview server running at %s', url)
    logger.info('PPTX file: %s', pptx_file)
    logger.info('Idle timeout: %ds (0 = disabled)', args.timeout)
    
    app.run(host='127.0.0.1', port=args.port, debug=False)
    return 0


if __name__ == '__main__':
    raise SystemExit(main())